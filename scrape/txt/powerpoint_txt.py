'''
Created on Jun 8, 2016

@author: pedrom
'''
from os import remove
import re
import subprocess

from pptx import Presentation

from utils.my_utils import is_good_lin


def pptx_handler(pp_file):
    prs = Presentation(pp_file)
    all_txt = ""
    for slide in prs.slides:
        slideTxt = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slideTxt += run.text+" "
        slideTxt = slideTxt.strip()
        slideTxt = re.sub('  +', ' ', slideTxt)
        if is_good_lin(slideTxt):
            all_txt += slideTxt+"\n"
    return all_txt
            
def ppt_handler(pp_file):
    command = "java -jar C:\\Users\\pedrom\\workspace\\GoogleScraper\\tika-app-1.13.jar --encoding=utf8 -t " + "\"" + pp_file + "\"> temp_ppt.txt"
    p = subprocess.Popen(command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    result = ""
    for line in p.stdout.readlines():
        result += str(line)
    retval = p.wait()
    with open("temp_ppt.txt", encoding="utf8") as ppFile:
        txt = ppFile.read()
        lins = [s for s in txt.splitlines(True) if s.strip("\r\n")]
        finalTxt = ""
        for lin in lins:
            lin = lin.strip()
            if is_good_lin(lin):
                finalTxt += lin + "\n"
    remove("temp_ppt.txt")
    return finalTxt[:-1]
        
#ppt_handler("C:\\Users\\pedrom\\workspace\\GoogleScraper\\testepp.ppt")